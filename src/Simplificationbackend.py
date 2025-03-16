from fastapi import FastAPI, File, UploadFile, HTTPException
from pydantic import BaseModel
from transformers import pipeline
import pdfplumber
from pptx import Presentation
import shutil
import os
import docx
import pytesseract
from PIL import Image
import io

app = FastAPI()

# Load Hugging Face summarization model
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

class TextRequest(BaseModel):
    text: str

@app.post("/summarize")
async def summarize_text(request: TextRequest):
    summary = summarizer(request.text, max_length=100, min_length=30, do_sample=False)
    return {"summary": summary[0]["summary_text"]}

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    file_ext = file.filename.split(".")[-1].lower()
    temp_file_path = f"temp_uploaded.{file_ext}"

    # Save file temporarily
    try:
        with open(temp_file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error saving file: {str(e)}")

    try:
        if file_ext == "txt":
            with open(temp_file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif file_ext == "pdf":
            text = extract_text_from_pdf(temp_file_path)
        elif file_ext in ["ppt", "pptx"]:
            text = extract_text_from_ppt(temp_file_path)
        elif file_ext in ["doc", "docx"]:
            text = extract_text_from_docx(temp_file_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type")

        if not text.strip():
            raise HTTPException(status_code=400, detail="No readable text found in file")

        # Summarize extracted text
        summary = summarizer(text, max_length=100, min_length=30, do_sample=False)
        return {"summary": summary[0]["summary_text"]}
    
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing file: {str(e)}")
    
    finally:
        if os.path.exists(temp_file_path):
            os.remove(temp_file_path)  # Clean up temporary file


def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                extracted_text = page.extract_text()
                if extracted_text:
                    text += extracted_text + "\n"
                else:
                    # Extract text from images using OCR (Tesseract)
                    for image in page.images:
                        img = Image.open(io.BytesIO(image["stream"].get_data()))
                        text += pytesseract.image_to_string(img) + "\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PDF: {str(e)}")
    return text.strip() if text else "No readable text found."


def extract_text_from_ppt(ppt_path):
    text = ""
    try:
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):  # Check if shape contains text
                    text += shape.text + "\n"
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                text += slide.notes_slide.notes_text_frame.text + "\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading PowerPoint: {str(e)}")
    return text.strip() if text else "No readable text found."


def extract_text_from_docx(docx_path):
    text = ""
    try:
        doc = docx.Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + "\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading Word document: {str(e)}")
    return text.strip() if text else "No readable text found."
